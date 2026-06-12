"""
title: Auto (Smart Routing)
author: Nexus
description: One model that does it all. Routes simple questions to the fast
  model, analysis to the reasoning model, coding to the coder model, and images
  to the vision model. ALSO handles document generation itself (Word / Excel /
  PowerPoint, and PDF-table -> Excel) because Pipe functions can't use Open
  WebUI tools. Hybrid export: simple files use fast deterministic builders;
  complex ones (charts, pivots, styling) are written by the coding model and run
  sandboxed (Claude-style). Strips the reasoning model's <think> output.
version: 0.3
"""
import glob
import hashlib
import io
import json
import os
import re
import sqlite3
import sys
import time
import uuid

import httpx
from pydantic import BaseModel, Field

# --------------------------------------------------------------------------- #
# Intent detection
# --------------------------------------------------------------------------- #
_REASONING = re.compile(
    r"\b(analy[sz]|assess|evaluat|compar|implication|risk|threat|root[\s-]?cause|"
    r"troubleshoot|diagnos|derive|calculat|prove|trade[\s-]?off|pros and cons|"
    r"step[\s-]?by[\s-]?step|reason|rationale|justif|design|architect|recommend|"
    r"gap|mitigat|consequence|why\b|how (?:would|could|should))", re.IGNORECASE,
)
_CODING = re.compile(
    r"(```|\b(code|coding|function|class |method|script|debug|compile|syntax|"
    r"stack ?trace|traceback|regex|sql query|api call|refactor|python|javascript|"
    r"typescript|java\b|c\+\+|golang|rust|bash script|"
    r"write (?:a |an )?(?:program|script|function|class)|"
    r"fix (?:this |my )?(?:code|bug|error)))", re.IGNORECASE,
)
_EXPORT = re.compile(
    r"\b(convert|export|make|save|turn|generate|create|put|download|give me)\b"
    r".{0,40}\b(word|docx|excel|xlsx|spreadsheet|powerpoint|pptx|slides?|deck|"
    r"presentation|document)\b", re.IGNORECASE,
)


def _is_coding(q):
    return bool(_CODING.search(q or ""))


def _needs_reasoning(q):
    q = (q or "").strip()
    return bool(_REASONING.search(q)) or len(q) > 300 or q.count("?") >= 2


def _route(text, has_image):
    """
    Decide which model handles a message. Pure + side-effect-free so it can be
    unit-tested (see test_router.py). Returns one of: "vision", "coding",
    "reasoning", "fast". Precedence: image > code > analysis > fast.
    """
    if has_image:
        return "vision"
    if _is_coding(text):
        return "coding"
    if _needs_reasoning(text):
        return "reasoning"
    return "fast"


def _export_format(q):
    if not _EXPORT.search(q or ""):
        return None
    s = q.lower()
    if re.search(r"excel|xlsx|spreadsheet", s):
        return "excel"
    if re.search(r"powerpoint|pptx|slides?|deck|presentation", s):
        return "pptx"
    if re.search(r"word|docx|document", s):
        return "word"
    return None


# A request is "complex" when the deterministic templates can't express it —
# charts, pivots, styling, formulas, multi-sheet, "format it like X". These go
# to the code-interpreter path (coding model writes a script that builds it).
_COMPLEX_EXPORT = re.compile(
    r"\b(chart|graph|plot|pivot|conditional format(?:ting)?|colou?r|highlight|"
    r"formula[es]?|multiple sheets?|multi[\s-]?sheet|dashboard|styl(?:e|ed|ing)|"
    r"theme|gridlines?|merge cells?|freeze panes?|borders?|"
    r"bar chart|pie chart|line chart|histogram|scatter|"
    r"format(?:ted|ting)? (?:it )?(?:like|as|to look)|"
    r"with (?:a |an )?(?:chart|graph|formula|total|sum|average))\b",
    re.IGNORECASE,
)


def _is_complex_export(q):
    return bool(_COMPLEX_EXPORT.search(q or ""))


_CODEGEN_EXT = {"excel": "xlsx", "word": "docx", "pptx": "pptx"}
_CODEGEN_LIBS = {
    "excel": "openpyxl (and matplotlib, saved as a PNG then inserted, if a "
             "chart is requested)",
    "word": "python-docx (and matplotlib for any chart image)",
    "pptx": "python-pptx (and matplotlib for any chart image)",
}


def _extract_code(text):
    """Pull the python code out of the model's ```python ...``` block."""
    blocks = re.findall(r"```(?:python|py)?\s*(.*?)```", text or "", re.DOTALL)
    if blocks:
        return max(blocks, key=len).strip()
    return (text or "").strip()


# --------------------------------------------------------------------------- #
# Sandboxing the code-interpreter path
# --------------------------------------------------------------------------- #
# The codegen path RUNS model-written Python. Defense in depth:
#   1. static denylist  — reject obviously-unsafe scripts before they run
#   2. rlimits          — cap memory / CPU / output size (kills runaways)
#   3. strict sandbox   — optional: no network + filesystem confined to workdir
# Layers 1-2 are always on; layer 3 is opt-in via the CODE_SANDBOX valve.

# An export script only needs office libs + matplotlib + OUTPUT_PATH. Anything
# reaching for the network, the shell, or absolute filesystem paths is rejected.
# Heuristic (evadable) — it backstops the subprocess limits, not replaces them.
_DENY = re.compile(
    r"(^\s*(?:import|from)\s+(?:socket|subprocess|requests|urllib|http|ftplib|"
    r"smtplib|telnetlib|ctypes|pty|asyncio|multiprocessing|paramiko|aiohttp))|"
    r"\bos\.(?:system|popen|exec[lv]\w*|spawn\w*|fork|remove|unlink|rmdir|"
    r"removedirs|kill)\b|"
    r"\bsubprocess\.|\bshutil\.(?:rmtree|move|copy\w*)\b|"
    r"\b(?:eval|exec|compile|__import__)\s*\(|"
    r"\bopen\s*\(\s*[frbu]*['\"]/",  # open() on a string-literal absolute path
    re.IGNORECASE | re.MULTILINE,
)


def _scan_code(code):
    """Return a rejection reason if the script trips the denylist, else None."""
    m = _DENY.search(code or "")
    if not m:
        return None
    return (f"blocked unsafe operation '{m.group(0).strip()[:60]}'. The script "
            "may ONLY build the file with openpyxl / python-docx / python-pptx "
            "/ matplotlib and write to os.environ['OUTPUT_PATH'] — no network, "
            "shell, subprocess, or absolute file paths.")


def _rlimits(mem_mb, cpu_s, fsize_mb):
    """preexec_fn: cap address space / CPU seconds / output size (POSIX)."""
    import resource
    mem = int(mem_mb) * 1024 * 1024
    for res, val in (
        (resource.RLIMIT_AS, mem),
        (resource.RLIMIT_DATA, mem),
        (resource.RLIMIT_CPU, int(cpu_s)),
        (resource.RLIMIT_FSIZE, int(fsize_mb) * 1024 * 1024),
        (resource.RLIMIT_NPROC, 128),
    ):
        try:
            resource.setrlimit(res, (val, val))
        except Exception:  # noqa: BLE001
            pass


def _sandbox_prefix(workdir, mode):
    """
    Strict mode -> a command prefix that runs the child with NO network and the
    filesystem confined to `workdir` (system dirs read-only). Prefers bubblewrap
    (bwrap); falls back to bare `unshare -n` (network isolation only). Returns
    ([], note) when no sandbox tool is present. (nsjail/firejail also work if you
    prefer them — wire them in here.)
    """
    import shutil
    if mode != "strict":
        return [], ""
    if shutil.which("bwrap"):
        return ([
            "bwrap", "--unshare-all", "--die-with-parent", "--new-session",
            "--ro-bind", "/usr", "/usr",
            "--ro-bind-try", "/bin", "/bin",
            "--ro-bind-try", "/sbin", "/sbin",
            "--ro-bind-try", "/lib", "/lib",
            "--ro-bind-try", "/lib64", "/lib64",
            "--ro-bind-try", "/etc", "/etc",
            "--proc", "/proc", "--dev", "/dev", "--tmpfs", "/tmp",
            "--bind", workdir, workdir, "--chdir", workdir,
        ], "")
    if shutil.which("unshare"):
        return (["unshare", "-n"],
                "network isolated; filesystem NOT confined (install bubblewrap "
                "for full strict sandboxing)")
    return ([], "strict sandbox requested but no sandbox tool found "
                "(install bubblewrap); ran with rlimits + denylist only")


def _run_codegen(code, workdir, out_path, timeout, mem_mb=4096, sandbox="basic"):
    """
    Execute model-written Python under layered isolation: static denylist, a
    minimal env (no inherited tokens/paths), a hard wall-clock timeout, rlimits
    (memory/CPU/output), and — in strict mode — no network + a confined
    filesystem. Returns (bytes, err).
    """
    import subprocess

    bad = _scan_code(code)
    if bad:
        return None, bad
    script = os.path.join(workdir, "gen.py")
    with open(script, "w") as f:
        f.write(code)
    env = {
        "PATH": os.environ.get("PATH", "/usr/local/bin:/usr/bin:/bin"),
        "HOME": workdir,
        "MPLBACKEND": "Agg",
        "MPLCONFIGDIR": workdir,
        "OUTPUT_PATH": out_path,
        "PYTHONDONTWRITEBYTECODE": "1",
    }
    prefix, note = _sandbox_prefix(workdir, sandbox)
    # Use the SAME interpreter running the pipe so the export libs (openpyxl,
    # matplotlib, …) are importable; bare "python3" could be a different env.
    cmd = prefix + [sys.executable or "python3", script]
    preexec = (lambda: _rlimits(mem_mb, max(5, timeout), 200)) \
        if os.name == "posix" else None
    try:
        proc = subprocess.run(
            cmd, cwd=workdir, capture_output=True, text=True, timeout=timeout,
            env=env, preexec_fn=preexec,
        )
    except subprocess.TimeoutExpired:
        return None, f"Script exceeded the {timeout}s time limit."
    except Exception as e:  # noqa: BLE001
        return None, f"Could not run script: {e}"
    if os.path.exists(out_path) and os.path.getsize(out_path) > 0:
        with open(out_path, "rb") as f:
            return f.read(), ""
    err = (proc.stderr or proc.stdout or "No file was produced.").strip()[-1800:]
    return None, (f"[{note}] {err}" if note else err)


# --------------------------------------------------------------------------- #
# File helpers
# --------------------------------------------------------------------------- #
def _data_dir():
    return os.environ.get("DATA_DIR", "/workspace/open-webui")


def _safe_name(filename, default="file"):
    """Strip any path components / unsafe chars to prevent path traversal."""
    base = os.path.basename(filename or default)
    base = re.sub(r"[^A-Za-z0-9._-]", "_", base).strip("._") or default
    return base[:80]


def _latest_pdf():
    """Most recently uploaded PDF (the one just attached)."""
    db = os.path.join(_data_dir(), "webui.db")
    try:
        c = sqlite3.connect(db)
        row = c.execute(
            "SELECT path, filename FROM file WHERE filename LIKE '%.pdf' "
            "ORDER BY created_at DESC LIMIT 1").fetchone()
        if row and row[0] and os.path.exists(row[0]):
            return row[0], row[1]
    except Exception:  # noqa: BLE001
        pass
    up = os.path.join(_data_dir(), "uploads")
    pdfs = sorted(glob.glob(os.path.join(up, "*.pdf")), key=os.path.getmtime,
                  reverse=True)
    if pdfs:
        return pdfs[0], os.path.basename(pdfs[0])
    return None, None


def _save_and_link(data, filename, mime, user_id):
    """
    Save the generated file into Open WebUI's storage and return a markdown link
    to it. We use a real /api/v1/files URL (not a data: link) because Open WebUI
    strips data: links from rendered markdown.
    """
    fid = str(uuid.uuid4())
    filename = _safe_name(filename, "export")
    up = os.path.join(_data_dir(), "uploads")
    os.makedirs(up, exist_ok=True)
    path = os.path.join(up, f"{fid}_{filename}")
    with open(path, "wb") as f:
        f.write(data)
    if not user_id:  # fall back to the most recent uploader (single-user case)
        try:
            c = sqlite3.connect(os.path.join(_data_dir(), "webui.db"))
            row = c.execute("SELECT user_id FROM file ORDER BY created_at DESC "
                            "LIMIT 1").fetchone()
            user_id = row[0] if row else ""
        except Exception:  # noqa: BLE001
            user_id = ""
    meta = json.dumps({"name": filename, "content_type": mime,
                       "size": len(data), "data": {}})
    now = int(time.time())
    h = hashlib.sha256(data).hexdigest()
    c = sqlite3.connect(os.path.join(_data_dir(), "webui.db"))
    c.execute(
        "INSERT INTO file (id,user_id,filename,meta,created_at,hash,data,"
        "updated_at,path) VALUES (?,?,?,?,?,?,?,?,?)",
        (fid, user_id, filename, meta, now, h, "{}", now, path),
    )
    c.commit()
    return f"📎 **[Download {filename}](/api/v1/files/{fid}/content)**"


_XLSX = "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
_DOCX = ("application/vnd.openxmlformats-officedocument."
         "wordprocessingml.document")
_PPTX = ("application/vnd.openxmlformats-officedocument."
         "presentationml.presentation")


def _pdf_tables_to_xlsx(path, name):
    import pdfplumber
    from openpyxl import Workbook
    from openpyxl.styles import Font
    wb = Workbook()
    wb.remove(wb.active)
    n = 0
    with pdfplumber.open(path) as pdf:
        for pno, page in enumerate(pdf.pages, 1):
            for tno, tbl in enumerate(page.extract_tables() or [], 1):
                if not tbl:
                    continue
                n += 1
                ws = wb.create_sheet(title=f"p{pno}_t{tno}"[:31])
                for r, row in enumerate(tbl, 1):
                    for col, val in enumerate(row, 1):
                        cell = ws.cell(row=r, column=col,
                                       value=("" if val is None else str(val)))
                        if r == 1:
                            cell.font = Font(bold=True)
    if n == 0:
        return None, 0
    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue(), n


_VISION_TABLE_PROMPT = (
    "You are reading one page of a document image. Extract EVERY table you see, "
    "exactly as laid out. Output each table as pipe-separated rows: one row per "
    "line, cells separated by '|', with the header row first. Preserve all rows "
    "and columns; keep empty cells as empty. Separate multiple tables with a "
    "line containing only '---'. Do not add commentary. If there is no table at "
    "all, output exactly NONE."
)


def _vision_pdf_to_xlsx(path, model, ollama_url, dpi=220):
    """
    Read tables by RENDERING each page and having the vision model read them —
    the 'like Claude' approach. Far better on complex/messy tables than the
    geometric parser. Returns (xlsx_bytes, n_tables) or (None, 0) on failure.
    """
    try:
        import fitz  # PyMuPDF
        import ollama
        from openpyxl import Workbook
        from openpyxl.styles import Font
    except Exception:  # noqa: BLE001
        return None, 0
    client = ollama.Client(host=ollama_url)
    wb = Workbook()
    wb.remove(wb.active)
    n = 0
    try:
        with fitz.open(path) as doc:
            for pno, page in enumerate(doc, 1):
                pix = page.get_pixmap(dpi=dpi)
                resp = client.chat(
                    model=model, keep_alive=-1,
                    messages=[{"role": "user", "images": [pix.tobytes("png")],
                               "content": _VISION_TABLE_PROMPT}],
                )
                txt = (resp["message"]["content"] or "").strip()
                if txt.upper().startswith("NONE") and "|" not in txt:
                    continue
                for block in txt.split("\n---\n"):
                    rows = [ln.split("|") for ln in block.splitlines()
                            if "|" in ln]
                    if not rows:
                        continue
                    n += 1
                    ws = wb.create_sheet(title=f"p{pno}_t{n}"[:31])
                    for r, row in enumerate(rows, 1):
                        for c, val in enumerate(row, 1):
                            cell = ws.cell(row=r, column=c, value=val.strip())
                            if r == 1:
                                cell.font = Font(bold=True)
    except Exception:  # noqa: BLE001
        return None, 0
    if n == 0:
        return None, 0
    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue(), n


def _md_table_to_xlsx(text):
    """Pull the first markdown/pipe table out of text into an xlsx."""
    from openpyxl import Workbook
    from openpyxl.styles import Font
    rows = []
    for line in (text or "").splitlines():
        if "|" in line and not re.match(r"^\s*\|?[\s:|-]+\|?\s*$", line):
            cells = [c.strip() for c in line.strip().strip("|").split("|")]
            rows.append(cells)
    if not rows:
        return None
    wb = Workbook()
    ws = wb.active
    for r, row in enumerate(rows, 1):
        for c, val in enumerate(row, 1):
            cell = ws.cell(row=r, column=c, value=val)
            if r == 1:
                cell.font = Font(bold=True)
    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue()


def _text_to_docx(title, body):
    from docx import Document
    doc = Document()
    if title:
        doc.add_heading(title, 0)
    for line in (body or "").split("\n"):
        s = line.strip()
        if not s:
            continue
        if s.startswith("## "):
            doc.add_heading(s[3:], 2)
        elif s.startswith("# "):
            doc.add_heading(s[2:], 1)
        elif s.startswith(("- ", "* ")):
            doc.add_paragraph(s[2:], style="List Bullet")
        else:
            doc.add_paragraph(re.sub(r"[#*`]", "", s))
    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


def _text_to_pptx(title, body):
    from pptx import Presentation
    prs = Presentation()
    s0 = prs.slides.add_slide(prs.slide_layouts[0])
    s0.shapes.title.text = title or "Presentation"
    # Each heading starts a new slide; bullets fill it.
    cur = None
    for line in (body or "").split("\n"):
        s = line.strip()
        if not s:
            continue
        if s.startswith("#") or (not s.startswith(("- ", "* ")) and cur is None):
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = re.sub(r"[#*`]", "", s).strip()
            cur = slide.placeholders[1].text_frame
            cur.clear()
            cur._first = True
        elif cur is not None:
            t = re.sub(r"^[-*]\s*", "", s)
            p = cur.paragraphs[0] if getattr(cur, "_first", False) else cur.add_paragraph()
            p.text = re.sub(r"[#*`]", "", t)
            cur._first = False
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


# --------------------------------------------------------------------------- #
# Pipe
# --------------------------------------------------------------------------- #
class Pipe:
    class Valves(BaseModel):
        OLLAMA_URL: str = Field(default="http://localhost:11434")
        # Model names default from models.env (the single source of truth) when
        # boot.sh has sourced it into Open WebUI's environment; the literals are
        # the same values as a fallback. 35B = fast/daily; 122B = deep analysis
        # (needs an H200/141GB-class GPU); 27B = coding; VL 32B = vision.
        FAST_MODEL: str = Field(default=os.getenv("FAST_MODEL", "qwen3.5:35b"))
        REASONING_MODEL: str = Field(default=os.getenv("REASONING_MODEL", "qwen3.5:122b"))
        CODE_MODEL: str = Field(default=os.getenv("CODE_MODEL", "qwen3.6:27b"))
        VISION_MODEL: str = Field(default=os.getenv("VISION_MODEL", "qwen3-vl:32b"))
        # Table extraction: "vision" reads tables with TABLE_MODEL (accurate,
        # slower); "fast" uses the geometric parser (instant, less accurate).
        TABLE_MODEL: str = Field(default=os.getenv("TABLE_MODEL", "qwen3-vl:32b"))
        TABLE_EXTRACTION: str = Field(default="vision")
        # Hybrid export: when a file request is "complex" (charts, pivots,
        # styling, formulas), the coding model writes a script that builds it
        # (Claude-style), executed in a sandboxed subprocess. Simple exports and
        # PDF->Excel keep the fast, deterministic builders. Falls back to
        # deterministic if the generated code fails after one retry.
        CODE_INTERPRETER: bool = Field(default=True)
        CODE_TIMEOUT: int = Field(default=90)
        # Sandboxing for the codegen path (it runs model-written Python):
        #   "basic"  = denylist + rlimits (memory/CPU/output caps). Always safe,
        #              no extra packages. Good for a trusted single-user box.
        #   "strict" = additionally no network + filesystem confined to a temp
        #              dir (needs `bwrap`/bubblewrap installed; falls back to
        #              `unshare -n` for network-only). Use for shared/client
        #              deployments. See README "Security — code execution".
        CODE_SANDBOX: str = Field(default="basic")
        CODE_MEM_MB: int = Field(default=4096)

    def __init__(self):
        self.valves = self.Valves()

    def pipes(self):
        return [{"id": "auto-router", "name": "Auto (Smart Routing)"}]

    @staticmethod
    def _normalize(m):
        content = m.get("content", "")
        if not isinstance(content, list):
            return {"role": m.get("role"), "content": content}
        texts, imgs = [], []
        for part in content:
            if not isinstance(part, dict):
                continue
            if part.get("type") == "text":
                texts.append(part.get("text", ""))
            elif part.get("type") == "image_url":
                url = part.get("image_url", {}).get("url", "")
                if "base64," in url:
                    imgs.append(url.split("base64,", 1)[1])
        out = {"role": m.get("role"), "content": " ".join(texts)}
        if imgs:
            out["images"] = imgs
        return out

    def _codegen_export(self, fmt, request, source):
        """
        Claude-style path: ask the coding model for a complete script that
        builds the file, run it sandboxed, retry once with the traceback if it
        errors. Returns (bytes, error_string).
        """
        import shutil
        import tempfile

        ext = _CODEGEN_EXT[fmt]
        workdir = tempfile.mkdtemp(prefix="nexus_gen_")
        out_path = os.path.join(workdir, f"output.{ext}")
        sys_prompt = (
            "You are a Python code generator. Write ONE complete, self-contained "
            f"Python 3 script that builds a .{ext} file using "
            f"{_CODEGEN_LIBS[fmt]}. "
            "Read the destination path from the environment variable OUTPUT_PATH "
            "(os.environ['OUTPUT_PATH']) and save the finished file there. "
            "Use only those libraries plus the Python standard library. Do NOT "
            "access the network, read other files, or print the data. Output "
            "ONLY the code inside a single ```python code block."
        )
        user_prompt = (
            f"Request: {request}\n\n"
            "Source content to base the file on (may include markdown tables, "
            f"lists, or prose):\n\"\"\"\n{(source or '')[:8000]}\n\"\"\""
        )
        msgs = [{"role": "system", "content": sys_prompt},
                {"role": "user", "content": user_prompt}]
        last_err = ""
        try:
            for _ in range(2):
                try:
                    resp = httpx.post(
                        f"{self.valves.OLLAMA_URL}/api/chat",
                        json={"model": self.valves.CODE_MODEL, "messages": msgs,
                              "stream": False, "keep_alive": -1},
                        timeout=300,
                    ).json()
                    raw = resp.get("message", {}).get("content", "")
                except Exception as e:  # noqa: BLE001
                    return None, f"coding model error: {e}"
                code = _extract_code(raw)
                if os.path.exists(out_path):
                    try:
                        os.remove(out_path)
                    except Exception:  # noqa: BLE001
                        pass
                data, err = _run_codegen(
                    code, workdir, out_path, self.valves.CODE_TIMEOUT,
                    self.valves.CODE_MEM_MB, self.valves.CODE_SANDBOX)
                if data:
                    return data, ""
                last_err = err
                msgs += [
                    {"role": "assistant", "content": raw},
                    {"role": "user", "content":
                        f"That script failed:\n{err}\n\nFix it and output only "
                        "the corrected complete script."},
                ]
            return None, last_err
        finally:
            shutil.rmtree(workdir, ignore_errors=True)

    def _handle_export(self, fmt, last_user, messages, user_id):
        """Generate a file directly and return a download-link string."""
        # Content to export = the previous assistant message (what the user is
        # usually referring to), else the user's own text.
        prev = ""
        for m in reversed(messages):
            if m.get("role") == "assistant" and m.get("content"):
                prev = m["content"]
                break
        title = "Export"

        # Hybrid: complex/custom requests -> coding model writes the file.
        fallback_note = ""
        if self.valves.CODE_INTERPRETER and _is_complex_export(last_user):
            mime = {"excel": _XLSX, "word": _DOCX, "pptx": _PPTX}[fmt]
            try:
                data, err = self._codegen_export(fmt, last_user, prev or last_user)
            except Exception as e:  # noqa: BLE001
                data, err = None, str(e)
            if data:
                fname = f"export.{_CODEGEN_EXT[fmt]}"
                return (f"Built it with the coding model ({self.valves.CODE_MODEL}).\n\n"
                        + _save_and_link(data, fname, mime, user_id))
            fallback_note = ("\n\n_(The custom-code build didn't produce a valid "
                             "file, so I used the standard builder.)_")

        if fmt == "excel":
            pdf_path, pdf_name = _latest_pdf()
            if pdf_path:
                data, n, how = None, 0, ""
                # Accuracy-first: read tables visually, fall back to geometric.
                if self.valves.TABLE_EXTRACTION == "vision":
                    data, n = _vision_pdf_to_xlsx(
                        pdf_path, self.valves.TABLE_MODEL, self.valves.OLLAMA_URL)
                    how = f" (read visually with {self.valves.TABLE_MODEL})"
                if not data:
                    data, n = _pdf_tables_to_xlsx(pdf_path, pdf_name)
                    how = ""
                if not data:  # last resort: small vision model
                    data, n = _vision_pdf_to_xlsx(
                        pdf_path, self.valves.VISION_MODEL, self.valves.OLLAMA_URL)
                    how = f" (read with {self.valves.VISION_MODEL})"
                if data:
                    fname = pdf_name.rsplit(".", 1)[0] + ".xlsx"
                    return (f"Extracted **{n} table(s)** from {pdf_name}{how}.\n\n"
                            + _save_and_link(data, fname, _XLSX, user_id)
                            + fallback_note)
            data = _md_table_to_xlsx(prev)
            if data:
                return ("Exported the table to Excel.\n\n"
                        + _save_and_link(data, "export.xlsx", _XLSX, user_id)
                        + fallback_note)
            return ("I couldn't find a table to put in Excel. Attach a PDF with "
                    "tables, or ask me to produce the data as a table first."
                    + fallback_note)

        if fmt == "word":
            data = _text_to_docx(title, prev or last_user)
            return ("Exported to Word.\n\n"
                    + _save_and_link(data, "document.docx", _DOCX, user_id)
                    + fallback_note)

        if fmt == "pptx":
            data = _text_to_pptx(title, prev or last_user)
            return ("Exported to PowerPoint.\n\n"
                    + _save_and_link(data, "presentation.pptx", _PPTX, user_id)
                    + fallback_note)
        return "Unsupported export format."

    def pipe(self, body, __user__=None):
        messages = [self._normalize(m) for m in body.get("messages", [])]
        last_user, last_has_image = "", False
        for m in reversed(messages):
            if m.get("role") == "user":
                last_user = m.get("content", "")
                last_has_image = bool(m.get("images"))
                break

        # 1) Export / convert intent -> generate the file ourselves.
        fmt = _export_format(last_user)
        if fmt:
            user_id = (__user__ or {}).get("id") if isinstance(__user__, dict) else None
            try:
                yield self._handle_export(fmt, last_user, messages, user_id)
            except Exception as e:  # noqa: BLE001
                yield f"⚠️ Export failed: {e}"
            return

        # 2) Route the chat. The Qwen3 models are "thinking" models: left on,
        # they emit a long hidden chain-of-thought BEFORE any answer, which
        # looks like a frozen UI (no content streams during thinking). So we
        # enable thinking ONLY for the deep-reasoning route; fast / vision /
        # coding answer directly for snappy responses.
        model, think = {
            "vision": (self.valves.VISION_MODEL, False),
            "coding": (self.valves.CODE_MODEL, False),
            "reasoning": (self.valves.REASONING_MODEL, True),
            "fast": (self.valves.FAST_MODEL, False),
        }[_route(last_user, last_has_image)]

        # The Qwen3 models stream reasoning in a separate `thinking` field
        # (not `content`). If we ignored it, the UI would sit silent for the
        # whole think phase, then dump the answer. Instead we re-wrap thinking
        # in <think></think> tags as it streams: Open WebUI renders that as a
        # live, collapsible "Thinking" section — real-time tokens, clean final
        # answer. Works whether thinking is long or short.
        think_open = False
        answer_started = False
        with httpx.stream(
            "POST", f"{self.valves.OLLAMA_URL}/api/chat",
            json={"model": model, "messages": messages, "think": think,
                  "stream": True, "keep_alive": -1},
            timeout=600,
        ) as r:
            for line in r.iter_lines():
                if not line:
                    continue
                try:
                    data = json.loads(line)
                except Exception:  # noqa: BLE001
                    continue
                msg = data.get("message", {}) or {}
                t = msg.get("thinking") or ""
                c = msg.get("content") or ""
                if t and not answer_started:
                    if not think_open:
                        think_open = True
                        yield "<think>"
                    yield t
                if c:
                    if think_open and not answer_started:
                        yield "</think>\n\n"
                    answer_started = True
                    yield c
        if think_open and not answer_started:
            yield "</think>"
